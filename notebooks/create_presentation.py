#!/usr/bin/env python3
"""
Script para generar la presentación del proyecto RL Meta-Allocator.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

# Alias for convenience
RgbColor = RGBColor


def add_title_slide(prs, title, subtitle=""):
    """Añade slide de título."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(0, 51, 102)  # Dark blue
    background.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title, number=""):
    """Añade slide de sección."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Banda lateral
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RgbColor(0, 102, 204)
    band.line.fill.background()

    # Número de sección
    if number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(1.5), Inches(1))
        tf = num_box.text_frame
        tf.paragraphs[0].text = number
        tf.paragraphs[0].font.size = Pt(72)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RgbColor(0, 102, 204)

    # Título
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(7), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(50, 50, 50)

    return slide


def add_content_slide(prs, title, bullets, subbullets=None):
    """Añade slide con contenido bullet."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Banda superior
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RgbColor(0, 51, 102)
    band.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)

    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    subbullets = subbullets or {}

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(12)

        # Subbullets
        if i in subbullets:
            for sub in subbullets[i]:
                p = tf.add_paragraph()
                p.text = f"    ‣ {sub}"
                p.font.size = Pt(18)
                p.font.color.rgb = RgbColor(100, 100, 100)
                p.space_after = Pt(6)

    return slide


def add_diagram_slide(prs, title, content_text):
    """Añade slide con diagrama/código."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Banda superior
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RgbColor(0, 51, 102)
    band.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)

    # Caja de contenido
    content_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    )
    content_shape.fill.solid()
    content_shape.fill.fore_color.rgb = RgbColor(245, 245, 245)
    content_shape.line.color.rgb = RgbColor(200, 200, 200)

    # Texto
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(8.6), Inches(4.6))
    tf = content_box.text_frame
    tf.word_wrap = True

    lines = content_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.name = "Consolas"
        p.font.color.rgb = RgbColor(50, 50, 50)

    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Añade slide con dos columnas."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Banda superior
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RgbColor(0, 51, 102)
    band.line.fill.background()

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)

    # Columna izquierda - título
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    tf.paragraphs[0].text = left_title
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(0, 102, 204)

    # Columna izquierda - contenido
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(8)

    # Columna derecha - título
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    tf.paragraphs[0].text = right_title
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(0, 102, 204)

    # Columna derecha - contenido
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(50, 50, 50)
        p.space_after = Pt(8)

    return slide


def add_table_slide(prs, title, headers, rows):
    """Añade slide con tabla."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Banda superior
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RgbColor(0, 51, 102)
    band.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)

    # Tabla
    n_rows = len(rows) + 1
    n_cols = len(headers)

    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(0.4 * n_rows)
    ).table

    # Headers
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RgbColor(0, 102, 204)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # Rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(240, 240, 240)

    return slide


def create_presentation():
    """Crea la presentación completa."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # PORTADA
    # =========================================================================
    add_title_slide(
        prs,
        "RL Meta-Allocator",
        "Asignación Dinámica de Capital mediante\nAprendizaje por Refuerzo"
    )

    # =========================================================================
    # AGENDA
    # =========================================================================
    add_content_slide(prs, "Agenda", [
        "Contexto y Motivación del Problema",
        "Análisis del Benchmark y Algoritmos",
        "Baselines Clásicos (MPT)",
        "Entorno RL Custom",
        "Agentes de Reinforcement Learning",
        "Resultados y Evaluación",
        "Conclusiones y Trabajo Futuro",
    ])

    # =========================================================================
    # SECCIÓN 1: CONTEXTO
    # =========================================================================
    add_section_slide(prs, "Contexto y Motivación", "01")

    add_content_slide(prs, "El Problema", [
        "Recibimos transacciones de N algoritmos caja negra",
        "No sabemos qué hacen internamente los algoritmos",
        "Tenemos un benchmark que invierte en estos algoritmos",
        "Objetivo: diseñar un meta-allocator que supere al benchmark",
    ], {
        0: ["Solo vemos entradas/salidas, no la lógica interna"],
        2: ["Conocemos sus trades y performance histórica"],
        3: ["Asignar capital dinámicamente entre algoritmos"],
    })

    add_content_slide(prs, "Restricción Crítica: Comparar Peras con Peras", [
        "El benchmark es la REFERENCIA, no el objetivo a replicar",
        "Debemos analizar sus características operativas:",
        "Restricciones y fricciones EQUIVALENTES para comparación justa",
        "El meta-allocator puede diferir en estrategia, no en condiciones",
    ], {
        1: ["Frecuencia de rebalanceo", "Duración de posiciones",
            "Concentración (HHI)", "Perfil de riesgo"],
    })

    add_content_slide(prs, "Enfoque: RL como Meta-Allocator", [
        "Reinforcement Learning para decisiones secuenciales de asignación",
        "El agente observa el estado del mercado y algoritmos",
        "Decide los pesos óptimos del portfolio",
        "Aprende a maximizar retornos ajustados por riesgo",
    ], {
        0: ["No predice precios, asigna capital"],
        2: ["Sujeto a restricciones realistas (turnover, concentración)"],
    })

    add_diagram_slide(prs, "Arquitectura del Sistema", """
    ┌─────────────────────────────────────────────────────────────────┐
    │                        DATOS DE ENTRADA                         │
    │  Algoritmos (OHLC)  │  Benchmark (trades)  │  Benchmarks macro  │
    └──────────────────────────────┬──────────────────────────────────┘
                                   │
                                   ▼
    ┌─────────────────────────────────────────────────────────────────┐
    │                    PHASE 1-2: ANÁLISIS                          │
    │  Feature Engineering  │  Clustering  │  Regime Inference        │
    └──────────────────────────────┬──────────────────────────────────┘
                                   │
                    ┌──────────────┴──────────────┐
                    ▼                             ▼
    ┌───────────────────────────┐   ┌───────────────────────────────┐
    │   PHASE 3: BASELINES      │   │   PHASE 4-5: RL AGENTS        │
    │  EW, RP, MinVar, MaxSharpe│   │  PPO, SAC, TD3                │
    │  Momentum, Vol Targeting  │   │  Entorno Custom (no Gym)      │
    └───────────────┬───────────┘   └───────────────┬───────────────┘
                    │                               │
                    └───────────────┬───────────────┘
                                    ▼
    ┌─────────────────────────────────────────────────────────────────┐
    │                    PHASE 6: EVALUACIÓN                          │
    │  Walk-Forward  │  Métricas comparables  │  Análisis de riesgo   │
    └─────────────────────────────────────────────────────────────────┘
    """)

    # =========================================================================
    # SECCIÓN 2: ANÁLISIS
    # =========================================================================
    add_section_slide(prs, "Análisis del Benchmark y Algoritmos", "02")

    add_content_slide(prs, "Phase 1: Data Pipeline", [
        "Carga y preprocesamiento de datos",
        "Detección automática de formatos (fechas, OHLC)",
        "Reconstrucción de equity curves y retornos",
        "Feature engineering (cumulative + rolling + cross-sectional)",
        "Inferencia de activo subyacente (two-stage)",
    ], {
        3: ["Ventanas: 5, 21, 63, 126, 252 días"],
        4: ["Fast Pearson screen → Deep 6-signal analysis"],
    })

    add_content_slide(prs, "Phase 2: Reverse Engineering del Benchmark", [
        "Profiling de todos los algoritmos",
        "Análisis de la política implícita del benchmark:",
        "Inferencia de régimen latente (HMM)",
        "Clustering temporal de algoritmos",
    ], {
        1: ["Sizing policy", "Temporal patterns", "Risk management", "Regime behavior"],
        2: ["Capa A: Régimen de mercado global", "Capa B: Comportamiento del benchmark"],
    })

    add_two_column_slide(
        prs, "Características Extraídas del Benchmark",
        "Métricas Operativas",
        [
            "Frecuencia de rebalanceo",
            "Turnover anualizado",
            "Holding period promedio",
            "Concentración (HHI)",
            "Número de posiciones activas",
        ],
        "Métricas de Riesgo",
        [
            "Volatilidad anualizada",
            "Max Drawdown",
            "VaR 95%",
            "CVaR 95%",
            "Sharpe Ratio",
        ]
    )

    add_content_slide(prs, "Clustering de Algoritmos", [
        "Temporal clustering: evolución semanal de clusters",
        "Behavioral clustering: familias por comportamiento",
        "Correlation clustering: grupos por correlación",
        "Métodos: KMeans, GMM, Hierarchical, DBSCAN, HDBSCAN",
    ], {
        0: ["Tres horizontes: cumulative, weekly, monthly"],
        1: ["Features: volatilidad, drawdown, autocorrelación, etc."],
    })

    # =========================================================================
    # SECCIÓN 3: BASELINES
    # =========================================================================
    add_section_slide(prs, "Baselines Clásicos (MPT)", "03")

    add_table_slide(prs, "6 Estrategias Baseline",
        ["Estrategia", "Descripción", "Objetivo"],
        [
            ["Equal Weight", "Pesos iguales para todos", "Baseline naive"],
            ["Risk Parity", "Pesos inversamente proporcionales a volatilidad", "Igualar contribución al riesgo"],
            ["Min Variance", "Optimización para mínima varianza", "Minimizar riesgo total"],
            ["Max Sharpe", "Optimización para máximo Sharpe", "Mejor risk-adjusted return"],
            ["Momentum", "Sobreponderar winners recientes", "Capturar tendencias"],
            ["Vol Targeting", "Ajustar exposure por volatilidad", "Control de riesgo dinámico"],
        ]
    )

    add_content_slide(prs, "Restricciones Aplicadas (Calibradas al Benchmark)", [
        "max_weight_per_algo: 0.40 (40% máximo por algoritmo)",
        "min_weight_per_algo: 0.00 (permite no invertir)",
        "max_turnover_per_rebalance: 0.30 (30% máximo)",
        "max_total_exposure: 1.0 (fully invested)",
        "Rebalance frequency: semanal",
    ])

    add_content_slide(prs, "Fast Backtester Optimizado", [
        "Implementación con Numba JIT compilation",
        "Caching de matrices de covarianza",
        "Procesamiento vectorizado",
        "351 trials ejecutados (6 estrategias × múltiples configs)",
    ], {
        0: ["@numba.njit para cálculos críticos"],
        2: ["Sin loops Python sobre filas/fechas"],
    })

    # =========================================================================
    # SECCIÓN 4: ENTORNO RL
    # =========================================================================
    add_section_slide(prs, "Entorno RL Custom", "04")

    add_content_slide(prs, "¿Por qué NO usar Gymnasium?", [
        "Gymnasium está diseñado para simulaciones genéricas",
        "Necesitamos un simulador financiero realista:",
        "Solo usamos gymnasium.spaces para compatibilidad con SB3",
        "Diseño event-driven con costes y restricciones reales",
    ], {
        1: ["Costes de transacción realistas", "Restricciones de portfolio",
            "Datos históricos reales (no simulados)", "Fricciones de mercado"],
    })

    add_diagram_slide(prs, "Arquitectura del Entorno", """
    ┌─────────────────────────────────────────────────────────────────┐
    │                      TradingEnvironment                         │
    │                                                                 │
    │  ┌─────────────────┐    ┌─────────────────┐    ┌─────────────┐ │
    │  │ MarketSimulator │    │   CostModel     │    │ Constraints │ │
    │  │                 │    │                 │    │             │ │
    │  │ - algo returns  │    │ - fixed cost    │    │ - max weight│ │
    │  │ - event-driven  │    │ - spread (5bp)  │    │ - turnover  │ │
    │  │ - realistic     │    │ - slippage (2bp)│    │ - exposure  │ │
    │  │                 │    │ - market impact │    │             │ │
    │  └────────┬────────┘    └────────┬────────┘    └──────┬──────┘ │
    │           │                      │                    │        │
    │           └──────────────────────┼────────────────────┘        │
    │                                  │                             │
    │                                  ▼                             │
    │                        ┌─────────────────┐                     │
    │                        │ RewardFunction  │                     │
    │                        │                 │                     │
    │                        │ - pure_returns  │                     │
    │                        │ - alpha vs bench│                     │
    │                        │ - risk penalties│                     │
    │                        └─────────────────┘                     │
    └─────────────────────────────────────────────────────────────────┘
    """)

    add_content_slide(prs, "Modelo de Costes", [
        "Fixed cost per trade: 0 (simplificación)",
        "Spread: 5 basis points",
        "Slippage: 2 basis points",
        "Market impact: 0.1 × sqrt(trade_size)",
        "Conservador para evitar sobreestimar performance",
    ])

    add_content_slide(prs, "Función de Reward", [
        "Múltiples tipos implementados:",
        "PURE_RETURNS: Solo maximiza retornos (sin penalizaciones)",
        "ALPHA_PENALIZED: Exceso vs benchmark - penalizaciones",
        "RISK_ADJUSTED: Sharpe diferencial",
        "Penalizaciones: costes, turnover, drawdown, riesgo excesivo",
    ], {
        1: ["Útil para diagnóstico y experimentación"],
        2: ["Penaliza costes, turnover, drawdown, tracking error"],
    })

    # =========================================================================
    # SECCIÓN 5: AGENTES RL
    # =========================================================================
    add_section_slide(prs, "Agentes de Reinforcement Learning", "05")

    add_two_column_slide(
        prs, "Algoritmos Implementados",
        "PPO (Baseline Principal)",
        [
            "Proximal Policy Optimization",
            "On-policy, estable",
            "Clip range: 0.2",
            "Red MLP [256, 256]",
            "Bueno para acciones continuas",
        ],
        "SAC / TD3 (Alternativas)",
        [
            "Off-policy, sample efficient",
            "SAC: máxima entropía",
            "TD3: twin critics, delayed update",
            "Mejor para exploration",
            "Útil si reward ruidoso",
        ]
    )

    add_content_slide(prs, "Hiperparámetros PPO", [
        "learning_rate: 3e-4",
        "n_steps: 2048 (steps antes de update)",
        "batch_size: 64",
        "n_epochs: 10 (epochs por update)",
        "gamma: 0.99 (discount factor)",
        "gae_lambda: 0.95 (GAE)",
        "ent_coef: 0.01 (exploration)",
    ])

    add_content_slide(prs, "Espacio de Observación", [
        "Retornos de algoritmos (rolling features)",
        "Volatilidades y correlaciones",
        "Drawdowns actuales",
        "Pesos actuales del portfolio",
        "Features de régimen (si disponibles)",
        "Normalización con running statistics",
    ])

    add_content_slide(prs, "Espacio de Acción", [
        "Continuo: Box(0, 1, shape=(n_algorithms,))",
        "Representa pesos target del portfolio",
        "Se normaliza para sumar 1.0",
        "Constraints se aplican post-acción",
    ])

    # =========================================================================
    # SECCIÓN 6: EVALUACIÓN
    # =========================================================================
    add_section_slide(prs, "Resultados y Evaluación", "06")

    add_content_slide(prs, "Metodología Walk-Forward", [
        "Train window: 252 días (1 año)",
        "Validation window: 63 días (1 trimestre)",
        "Test window: 63 días (1 trimestre)",
        "Step size: 63 días (rolling)",
        "NUNCA train/test aleatorio - siempre temporal",
    ], {
        4: ["Evita data leakage", "Simula condiciones reales de trading"],
    })

    add_table_slide(prs, "Métricas de Evaluación",
        ["Categoría", "Métricas"],
        [
            ["Performance Absoluta", "Return, Volatility, Sharpe, Sortino, Calmar"],
            ["Performance Relativa", "Excess Return, Tracking Error, Information Ratio, Alpha, Beta"],
            ["Operativa", "Turnover, Holding Period, Concentration (HHI)"],
            ["Riesgo", "Max Drawdown, VaR 95%, CVaR 95%"],
        ]
    )

    add_content_slide(prs, "Criterio de Éxito", [
        "Batir al benchmark en métricas ajustadas por riesgo",
        "Superar a los baselines clásicos",
        "A igualdad de restricciones y fricciones",
        "Consistencia en walk-forward (no solo in-sample)",
    ])

    add_content_slide(prs, "Sanity Checks Implementados", [
        "EW constante = Baseline Equal Weight (verificar equivalencia)",
        "Costes = 0 → debería rotar más",
        "Turnover penalty alta → debería tender a buy-and-hold",
        "Replicar benchmark → reward ≈ 0",
        "Si falla alguno, hay bug en el código",
    ])

    # =========================================================================
    # SECCIÓN 7: CONCLUSIONES
    # =========================================================================
    add_section_slide(prs, "Conclusiones y Trabajo Futuro", "07")

    add_content_slide(prs, "Logros del Proyecto", [
        "Pipeline completo de análisis de datos financieros",
        "Reverse engineering del benchmark (características operativas)",
        "6 baselines clásicos implementados y evaluados",
        "Entorno RL custom con costes y restricciones realistas",
        "3 agentes RL (PPO, SAC, TD3) con configuración optimizada",
        "Framework de evaluación walk-forward riguroso",
    ])

    add_content_slide(prs, "Lecciones Aprendidas", [
        "RL para allocation es DIFÍCIL - señal débil, mucho ruido",
        "Los baselines clásicos son difíciles de superar",
        "La función de reward es CRÍTICA",
        "Simplificar primero (pure returns), luego añadir complejidad",
        "Walk-forward es esencial - in-sample engaña",
    ])

    add_content_slide(prs, "Trabajo Futuro", [
        "Enfoque híbrido: MPT base + RL tilts",
        "Offline RL con datos históricos (CQL/IQL)",
        "Más features de régimen (macro indicators)",
        "Ensemble de agentes",
        "Transfer learning entre regímenes",
    ], {
        0: ["target = risk_parity + bounded_rl_adjustment"],
    })

    add_content_slide(prs, "Stack Tecnológico", [
        "Python 3.13 + NumPy + Pandas",
        "Numba para JIT compilation",
        "Stable-Baselines3 para RL",
        "PyTorch como backend",
        "Scikit-learn para clustering/PCA",
        "HMMlearn para inferencia de regímenes",
    ])

    # =========================================================================
    # CIERRE
    # =========================================================================
    add_title_slide(
        prs,
        "¿Preguntas?",
        "Gracias por su atención"
    )

    # Guardar
    output_path = Path(__file__).parent / "RL_Meta_Allocator_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentación guardada en: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
